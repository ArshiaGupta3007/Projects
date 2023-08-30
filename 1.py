from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Create the title slide
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Man-made Disasters"
subtitle.text = "An Overview"

# Add content slides about different man-made disasters
disasters = [
    {
        "title": "Oil Spill",
        "description": "An oil spill is the release of a liquid petroleum hydrocarbon into the environment, " \
                       "especially the marine ecosystem, due to human activity. It causes significant damage " \
                       "to marine life, habitats, and the overall ecosystem.",
        "image_path": "oil_spill.jpg"  # Path to an image related to the disaster
    },
    {
        "title": "Nuclear Disaster",
        "description": "A nuclear disaster is a catastrophic failure of a nuclear power plant, leading to " \
                       "the release of radioactive materials into the environment. It poses severe health " \
                       "risks to humans and has long-term environmental impacts.",
        "image_path": "nuclear_disaster.jpg"
    },
    {
        "title": "Industrial Accident",
        "description": "An industrial accident refers to a sudden and unintended occurrence in an industrial " \
                       "setting that causes harm to people, property, or the environment. Examples include " \
                       "chemical spills, explosions, and fires.",
        "image_path": "industrial_accident.jpg"
    }
]

for disaster in disasters:
    slide_layout = presentation.slide_layouts[1]  # Content slide layout
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = disaster["title"]
    content.text = disaster["description"]

    # Add an image related to the disaster
    image_path = disaster["image_path"]
    if image_path:
        left = Inches(2)
        top = Inches(2)
        width = Inches(5)
        height = Inches(3)
        slide.shapes.add_picture(image_path, left, top, width, height)

# Save the presentation
presentation.save("man_made_disasters.pptx")